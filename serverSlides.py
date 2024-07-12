from flask import Flask, request, jsonify, send_file, send_from_directory
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import os

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static'

@app.route('/process', methods=['POST'])
def process_data():
    data = request.get_json()
    print(data)
    if not isinstance(data['data'], list):
        return jsonify({'error': 'Invalid input data'}), 400

    # Create a Presentation object
    prs = Presentation()

    # Iterate through the list and create slides
    for item in data['data']:
        if 'title' in item and 'content' in item:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use the title and content layout
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = item['title']
            content.text = item['content']
        else:
            return jsonify({'error': 'Invalid input data'}), 400

    # Save the presentation to a file
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    pptx_filename = f'presentation_{timestamp}.pptx'
    pptx_filepath = os.path.join(app.config['UPLOAD_FOLDER'], pptx_filename)
    prs.save(pptx_filepath)

    return send_file(pptx_filepath, as_attachment=True)
@app.route('/static/<path:filename>', methods=['GET'])
def download_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/files', methods=['GET'])
def list_files():
    files = os.listdir(app.config['UPLOAD_FOLDER'])
    return jsonify(files)
if __name__ == '__main__':
    app.run(debug=True, port=5000)
